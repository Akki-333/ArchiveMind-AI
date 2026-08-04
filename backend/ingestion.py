import io
import uuid
import time
import json
from fastapi import APIRouter, UploadFile, File, HTTPException, Depends
from auth import get_current_user
import PyPDF2
import docx
from pptx import Presentation
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_core.documents import Document
from langchain_core.prompts import ChatPromptTemplate
from langchain_core.output_parsers import JsonOutputParser
from langchain_pinecone import PineconeVectorStore
from database import pc, index_name, embeddings_model, llm, neo4j_driver

router = APIRouter()

# Graph Extraction Chain (Same as PoC but tailored for real DB)
parser = JsonOutputParser()
extraction_prompt = ChatPromptTemplate.from_messages([
    ("system", "You are an AI extracting entities and relationships from documents.\n"
               "Extract the main entities (such as concepts, technologies, frameworks, methods, organizations, or metrics) and how they relate.\n"
               "CRITICAL: You MUST ONLY extract entities and relationships that are explicitly mentioned in the provided text AND are relevant to the user's query.\n"
               "If the user's query is generic (e.g. 'hii', 'hello') or the text does not contain any clear entities relevant to the query, you MUST return empty lists.\n"
               "Return ONLY a valid JSON object with a 'nodes' list and an 'edges' list.\n"
               "Format:\n"
               "{{\n"
               "  \"nodes\": [ {{\"id\": \"Entity Name\", \"type\": \"Category\"}}, ...],\n"
               "  \"edges\": [ {{\"source\": \"Entity Name 1\", \"target\": \"Entity Name 2\", \"label\": \"relationship_type\"}}, ...]\n"
               "}}\n\n"
               "{format_instructions}"),
    ("human", "User Query: {query}\n\nExtract relevant entities from the following text:\n\n{text}")
])
extraction_chain = extraction_prompt | llm | parser

# Overview Profile Chain
overview_parser = JsonOutputParser()
overview_prompt = ChatPromptTemplate.from_messages([
    ("system", "You are an AI generating an overview for a government document.\n"
               "Read the provided text (the first few pages) and return ONLY a valid JSON object with a 'summary' string (max 3 sentences) and a 'key_entities' list of strings.\n"
               "Format:\n"
               "{{\n"
               "  \"summary\": \"Brief summary here...\",\n"
               "  \"key_entities\": [\"Entity1\", \"Entity2\"]\n"
               "}}\n\n"
               "{format_instructions}"),
    ("human", "Extract from the following text:\n\n{text}")
])
overview_chain = overview_prompt | llm | overview_parser

def save_to_neo4j(nodes, edges, doc_id):
    """Pushes extracted nodes and edges into the Neo4j AuraDB."""
    with neo4j_driver.session() as session:
        for node in nodes:
            # Merge ensures we don't create duplicates
            query = (
                "MATCH (d:Document {id: $doc_id}) "
                "MERGE (n:Entity {id: $id}) "
                "SET n.type = $type "
                "MERGE (n)-[:FOUND_IN]->(d)"
            )
            session.run(query, id=node["id"], type=node.get("type", "Unknown"), doc_id=doc_id)
            
        for edge in edges:
            # Create relationships between existing nodes
            query = (
                "MATCH (source:Entity {id: $source_id}) "
                "MATCH (target:Entity {id: $target_id}) "
                "MERGE (source)-[r:RELATED_TO {type: $label}]->(target)"
            )
            session.run(query, source_id=edge["source"], target_id=edge["target"], label=edge.get("label", "related_to"))

@router.post("/upload")
async def upload_document(file: UploadFile = File(...), username: str = Depends(get_current_user)):
    ext = file.filename.split('.')[-1].lower()
    allowed_exts = ["pdf", "docx", "pptx", "txt", "md", "csv"]
    if ext not in allowed_exts:
        raise HTTPException(status_code=400, detail=f"Unsupported file format. Allowed: {', '.join(allowed_exts)}")

    try:
        # Check upload limit
        with neo4j_driver.session() as session:
            count_result = session.run("MATCH (u:User {username: $username})-[:UPLOADED]->(d:Document) RETURN COUNT(d) AS count", username=username)
            if count_result.single()["count"] >= 5:
                raise HTTPException(status_code=400, detail="Upload limit reached. You can only upload a maximum of 5 documents.")
                
        doc_id = str(uuid.uuid4())
        ts = int(time.time() * 1000)
        # 1. Read File
        contents = await file.read()
        text = ""
        
        if ext == "pdf":
            pdf_reader = PyPDF2.PdfReader(io.BytesIO(contents))
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
        elif ext == "docx":
            doc = docx.Document(io.BytesIO(contents))
            for para in doc.paragraphs:
                text += para.text + "\n"
        elif ext == "pptx":
            prs = Presentation(io.BytesIO(contents))
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        elif ext in ["txt", "md", "csv"]:
            text = contents.decode('utf-8', errors='ignore')
            
        if not text.strip():
            raise HTTPException(status_code=400, detail="Could not extract text from document.")

        # 2. Chunking
        text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=100)
        chunks = text_splitter.split_text(text)
        
        # Convert to LangChain Documents, inject doc_id and filename
        documents = [Document(page_content=c, metadata={"source": file.filename, "doc_id": doc_id}) for c in chunks]

        # 3. Vector Brain: Push to Pinecone
        PineconeVectorStore.from_documents(
            documents, 
            embeddings_model, 
            index_name=index_name
        )

        # Generate Document Overview Profile from first 2 chunks
        overview_text = "\n".join(chunks[:2])
        try:
            overview_result = overview_chain.invoke({
                "text": overview_text,
                "format_instructions": overview_parser.get_format_instructions()
            })
            summary = overview_result.get("summary", "Summary not available.")
            key_entities = json.dumps(overview_result.get("key_entities", []))
        except Exception as e:
            print(f"Error extracting overview: {e}")
            summary = "Failed to generate summary."
            key_entities = "[]"
            
        # Save Document node to Neo4j
        with neo4j_driver.session() as session:
            session.run("""
                MATCH (u:User {username: $username})
                CREATE (u)-[:UPLOADED]->(d:Document {
                    id: $doc_id, 
                    filename: $filename, 
                    summary: $summary, 
                    key_entities: $key_entities, 
                    created_at: $ts
                })
            """, username=username, doc_id=doc_id, filename=file.filename, summary=summary, key_entities=key_entities, ts=ts)

        # 4. Relational Brain: Graph extraction is now deferred to query time!
        # We no longer extract the graph during upload, solving rate limits and speed issues.

        return {
            "status": "success", 
            "message": f"Document '{file.filename}' processed successfully.",
            "chunks_embedded": len(documents)
        }

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))
